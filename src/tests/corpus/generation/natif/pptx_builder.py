# ====== Code Summary ======
# PptxCorpusBuilder — builds a multilingual .pptx deck from a content pack: title slide, abstract,
# several bulleted content slides (cycling the section/paragraph pools for length), an embedded
# image slide, a native table slide and a native chart slide. Each slide becomes a PDF page after
# LibreOffice conversion. Driven by spec.doc_type + spec.language.

# ====== Standard Library Imports ======
from __future__ import annotations

import io

# ====== Third-Party Library Imports ======
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches, Pt

# ====== Local Project Imports ======
from .base import BaseDocumentBuilder
from .content import ContentPack, get_content
from .image_factory import ImageFactory

# Number of bulleted content slides (cycles the content pool).
_CONTENT_SLIDES = 6


class PptxCorpusBuilder(BaseDocumentBuilder):
    """Builds a multilingual slide deck for one (doc_type, language)."""

    def build(self) -> bytes:
        """
        Assemble the .pptx and return its bytes.

        Returns:
            bytes: A valid, multi-slide .pptx presentation.
        """
        # 1. Resolve content + title slide
        content = get_content(self.spec.doc_type, self.spec.language)
        prs = Presentation()
        self.__add_title_slide(prs, content)

        # 2. Bulleted content slides (cycled for length)
        for i in range(_CONTENT_SLIDES):
            self.__add_bullets_slide(prs, content, i)

        # 3. Image, native table, native chart slides
        self.__add_image_slide(prs, content)
        self.__add_table_slide(prs, content)
        self.__add_chart_slide(prs, content)

        # 4. Serialize
        buf = io.BytesIO()
        prs.save(buf)
        return buf.getvalue()

    def __add_title_slide(self, prs: Presentation, content: ContentPack) -> None:
        """Add the title slide (title + subtitle carrying the searchable phrase)."""
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = content.title
        slide.placeholders[1].text = f"{content.subtitle}\n{content.searchable_phrase}"

    @staticmethod
    def _blank(prs: Presentation, title: str) -> object:
        """Add a blank slide with a manual title textbox and return it."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
        frame = box.text_frame
        frame.text = title
        frame.paragraphs[0].font.size = Pt(28)
        frame.paragraphs[0].font.bold = True
        return slide

    def __add_bullets_slide(self, prs: Presentation, content: ContentPack, i: int) -> None:
        """Add a content slide: section title + a multi-level bulleted outline from the pool."""
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = content.section_title(i)
        body = slide.placeholders[1].text_frame
        body.text = content.para(i)
        for level, item in enumerate(content.list_items[:3]):
            para = body.add_paragraph()
            para.text = item
            para.level = min(level + 1, 4)

    def __add_image_slide(self, prs: Presentation, content: ContentPack) -> None:
        """Add a slide with a synthetic chart image."""
        slide = self._blank(prs, content.section_title(0))
        png = ImageFactory.bar_chart("KPI")
        slide.shapes.add_picture(io.BytesIO(png), Inches(1.5), Inches(1.5), width=Inches(6))

    def __add_table_slide(self, prs: Presentation, content: ContentPack) -> None:
        """Add a slide with a native table built from the content pack's table."""
        slide = self._blank(prs, content.table_caption)
        headers = content.table_headers
        rows = content.table_rows[:3]
        table = slide.shapes.add_table(
            len(rows) + 1, len(headers), Inches(0.6), Inches(1.6), Inches(8.8), Inches(3)
        ).table
        for c, head in enumerate(headers):
            table.cell(0, c).text = head
        for r, row in enumerate(rows, start=1):
            for c, value in enumerate(row):
                table.cell(r, c).text = str(value)

    def __add_chart_slide(self, prs: Presentation, content: ContentPack) -> None:
        """Add a slide with a native clustered-column chart from the table's numeric columns."""
        slide = self._blank(prs, content.section_title(1))
        chart_data = CategoryChartData()
        chart_data.categories = [row[0] for row in content.table_rows]
        # Use the second column (first metric) as the series, stripped of separators/symbols.
        series = [self.__num(row[1]) for row in content.table_rows]
        chart_data.add_series(content.table_headers[1], series)
        slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(1.5), Inches(1.6), Inches(7), Inches(4.5), chart_data
        )

    @staticmethod
    def __num(value: str) -> float:
        """Parse a numeric value from a localized string (strips spaces / thin spaces / commas)."""
        cleaned = value.replace(" ", "").replace(" ", "").replace(",", "").replace("%", "")
        try:
            return float(cleaned)
        except ValueError:
            return 0.0
